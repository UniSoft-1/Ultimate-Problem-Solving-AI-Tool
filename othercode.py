import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from sklearn.cluster import KMeans
from sklearn.preprocessing import StandardScaler
from sklearn.ensemble import IsolationForest
import plotly.express as px
import plotly.graph_objects as go
import streamlit as st
from sklearn.linear_model import LinearRegression
from scipy import stats
from scipy.stats import norm
from statsmodels.tsa.stattools import adfuller
from statsmodels.stats.anova import anova_lm
from statsmodels.formula.api import ols
from sklearn.decomposition import PCA
from statsmodels.tsa.holtwinters import ExponentialSmoothing  
from statsmodels.tsa.arima.model import ARIMA
from sklearn.metrics import mean_absolute_percentage_error, mean_absolute_error, mean_squared_error
from pptx import Presentation
from redx_tool import run_redx_tool
import scipy.stats as stats
from statsmodels.stats.weightstats import ttest_ind
import statsmodels.formula.api as smf

# Custom CSS for a Minimalistic Design
st.markdown("""
    <style>
        body {
            background-color: #f7f8fa; /* Light, clean background */
            font-family: 'Roboto', sans-serif;
            color: #333;
            margin: 0;
            padding: 0;
        }

        .menu-container {
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(300px, 1fr));
            gap: 24px;
            margin: 40px;
            padding: 0 20px;
        }

        .menu-card {
            background-color: #ffffff;
            padding: 25px;
            border-radius: 10px;
            box-shadow: 0 4px 12px rgba(0, 0, 0, 0.08);
            text-align: center;
            transition: transform 0.3s ease, box-shadow 0.3s ease;
            cursor: pointer;
            text-decoration: none;
            border: 1px solid #e1e1e1; /* Soft border */
        }

        .menu-card:hover {
            transform: translateY(-5px);
            box-shadow: 0 8px 20px rgba(0, 0, 0, 0.12);
        }

        .menu-card h3 {
            font-size: 1.4em;
            margin-bottom: 20px;
            color: #2d3748; /* Darker gray for the title */
        }

        .menu-card p {
            font-size: 1em;
            color: #4a5568; /* Soft gray text */
        }

        .menu-card img {
            width: 80px;
            height: 80px;
            margin-bottom: 20px;
            object-fit: cover;
            border-radius: 50%; /* Circular icons */
        }

        .stButton button {
            background-color: #4CAF50; /* Soft green */
            color: white;
            border-radius: 8px;
            padding: 12px 25px;
            font-size: 16px;
            cursor: pointer;
            border: none;
            box-shadow: 0 4px 8px rgba(0, 0, 0, 0.1);
            transition: background-color 0.3s, transform 0.3s;
        }

        .stButton button:hover {
            background-color: #45a049;
            transform: translateY(-2px);
        }

        .stButton button:active {
            transform: translateY(1px);
        }

        .sidebar .sidebar-content {
            background-color: #ffffff;
            padding: 10px;
            box-shadow: 0px 4px 8px rgba(0, 0, 0, 0.1);
            border-radius: 10px;
            box-sizing: border-box;
        }

        .sidebar .sidebar-content .element-container {
            margin-top: 20px;
        }

        h1, h2, h3 {
            font-weight: 500;
        }

        /* Responsive design for smaller screens */
        @media (max-width: 768px) {
            .menu-card {
                padding: 18px;
            }

            .menu-card h3 {
                font-size: 1.2em;
            }

            .menu-card p {
                font-size: 0.9em;
            }

            .menu-card img {
                width: 60px;
                height: 60px;
            }

            .stButton button {
                font-size: 14px;
            }
        }
    </style>
""", unsafe_allow_html=True)

import streamlit as st
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
import pandas as pd

def generate_8d_report(df, dmaic_insights=None, selected_analysis=None):
    st.write("### 8D Report Generation")

    # D1 - Team Formation
    team_members = st.text_input("Enter the team members involved in the problem-solving process:")
    if not team_members:
        st.warning("Please provide team members for D1.")
        return

    # D2 - Problem Description
    problem_description = st.text_area("Describe the problem you're trying to solve:")
    if not problem_description:
        st.warning("Please provide a description of the problem for D2.")
        return

    # D3 - Interim Containment Actions
    containment_actions = st.text_area("Describe the interim actions taken to contain the problem:")
    if not containment_actions:
        st.warning("Please provide interim containment actions for D3.")
        return

    # D4 - Root Cause Analysis
    root_cause = st.text_area("Based on the analysis, what is the root cause of the problem?")
    if not root_cause:
        st.warning("Please provide the root cause of the problem for D4.")
        return

    # D5 - Permanent Corrective Actions
    corrective_actions = st.text_area("What are the permanent corrective actions implemented?")
    if not corrective_actions:
        st.warning("Please provide permanent corrective actions for D5.")
        return

    # D6 - Implement Corrective Actions
    implementation_steps = st.text_area("Describe how the corrective actions were implemented:")
    if not implementation_steps:
        st.warning("Please describe the implementation steps for D6.")
        return

    # D7 - Prevent Recurrence
    prevention_measures = st.text_area("What measures are put in place to prevent recurrence?")
    if not prevention_measures:
        st.warning("Please provide preventive measures for D7.")
        return

    st.subheader("🔧 FMEA & Control Plan Adjustments")

    fmea_description = st.text_area("Describe your current FMEA setup (including severity, occurrence, detection ratings, etc.):")
    cp_description = st.text_area("Describe your current Control Plan (key parameters, controls, frequency):")

    fmea_file = st.file_uploader("Or upload your FMEA file (CSV or Excel)", type=["csv", "xlsx"])
    cp_file = st.file_uploader("Or upload your Control Plan file (CSV or Excel)", type=["csv", "xlsx"])

    fmea_changes = "FMEA analysis not provided."
    cp_changes = "Control Plan analysis not provided."

    # Basic logic to suggest changes
    if fmea_description:
        fmea_changes = f"Update severity or occurrence ratings based on the identified root cause: {root_cause}. Ensure detection methods align with new controls."
    elif fmea_file:
        try:
            fmea_df = pd.read_excel(fmea_file) if fmea_file.name.endswith(".xlsx") else pd.read_csv(fmea_file)
            fmea_changes = "Based on your FMEA, consider updating the rows related to the failure mode: '{}'. Review RPN scores and detection mechanisms.".format(root_cause)
        except Exception as e:
            fmea_changes = f"Error reading FMEA file: {e}"

    if cp_description:
        cp_changes = f"Ensure that the Control Plan now includes monitoring for the failure mode or root cause: {root_cause} with clear corrective checkpoints."
    elif cp_file:
        try:
            cp_df = pd.read_excel(cp_file) if cp_file.name.endswith(".xlsx") else pd.read_csv(cp_file)
            cp_changes = "Based on your Control Plan, include specific controls for the root cause and modify control frequency if needed."
        except Exception as e:
            cp_changes = f"Error reading Control Plan file: {e}"

    # D8 - Recognize the Team
    team_recognition = st.text_area("How will the team be recognized for their efforts?")
    if not team_recognition:
        st.warning("Please provide team recognition details for D8.")
        return

    st.write("### 📘 8D Report Summary")
    st.write(f"**D1**: {team_members}")
    st.write(f"**D2**: {problem_description}")
    st.write(f"**D3**: {containment_actions}")
    st.write(f"**D4**: {root_cause}")
    st.write(f"**D5**: {corrective_actions}")
    st.write(f"**D6**: {implementation_steps}")
    st.write(f"**D7**: {prevention_measures}")
    st.write("**FMEA Changes Suggested:**", fmea_changes)
    st.write("**Control Plan Changes Suggested:**", cp_changes)
    st.write(f"**D8**: {team_recognition}")

    # Generate chart
    fig, ax = plt.subplots(figsize=(10, 6))
    sns.histplot(df, ax=ax)
    ax.set_title("Problem Data Distribution")
    st.pyplot(fig)

    # Generate presentation
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "8D Report"
    content = slide.shapes.placeholders[1]
    content.text = (
        f"D1: {team_members}\nD2: {problem_description}\nD3: {containment_actions}\n"
        f"D4: {root_cause}\nD5: {corrective_actions}\nD6: {implementation_steps}\n"
        f"D7: {prevention_measures}\nFMEA Changes: {fmea_changes}\n"
        f"Control Plan Changes: {cp_changes}\nD8: {team_recognition}"
    )

    # Save presentation
    presentation_file = "/mnt/data/8d_report.pptx"
    prs.save(presentation_file)
    st.download_button("📥 Download 8D Report", presentation_file)

# Core Hypothesis Testing Tool
def hypothesis_testing_tool(df):
    st.subheader("Hypothesis Testing Tool")

    numeric_columns = df.select_dtypes(include=['number']).columns.tolist()
    categorical_columns = df.select_dtypes(exclude=['number']).columns.tolist()

    selected_variables = st.multiselect("Select Variable Columns for Analysis", numeric_columns)
    selected_category = st.selectbox("Select Categorical Column for Grouping (optional)", [None] + categorical_columns)

    if st.button("Run Hypothesis Tests"):
        if selected_category:
            categories = df[selected_category].dropna().unique()
            if len(categories) == 2:
                group1 = df[df[selected_category] == categories[0]][selected_variables]
                group2 = df[df[selected_category] == categories[1]][selected_variables]
                st.write(f"Comparing groups: {categories[0]} vs {categories[1]}")
                for col in selected_variables:
                    t_stat, p_val, _ = ttest_ind(group1[col].dropna(), group2[col].dropna())
                    st.write(f"T-test for {col}: t = {t_stat:.4f}, p = {p_val:.4f}")
                    interpret_pval(p_val)
            elif len(categories) > 2:
                st.info("One-Way ANOVA will be used for more than 2 categories.")
                for col in selected_variables:
                    groups = [df[df[selected_category] == cat][col].dropna() for cat in categories]
                    f_stat, p_val = stats.f_oneway(*groups)
                    st.write(f"ANOVA for {col} by {selected_category}: F = {f_stat:.4f}, p = {p_val:.4f}")
                    interpret_pval(p_val)
            else:
                st.warning("Chi-Square Test will be used.")
                for col in selected_variables:
                    cont_table = pd.crosstab(df[col], df[selected_category])
                    chi2, p, dof, expected = stats.chi2_contingency(cont_table)
                    st.write(f"Chi-square for {col} vs {selected_category}: chi2 = {chi2:.4f}, p = {p:.4f}")
                    interpret_pval(p)
        else:
            if len(selected_variables) >= 2:
                for i in range(len(selected_variables)):
                    for j in range(i+1, len(selected_variables)):
                        col1 = selected_variables[i]
                        col2 = selected_variables[j]
                        t_stat, p_val = stats.ttest_ind(df[col1].dropna(), df[col2].dropna())
                        st.write(f"T-test: {col1} vs {col2}: t = {t_stat:.4f}, p = {p_val:.4f}")
                        interpret_pval(p_val)
            else:
                st.error("Select at least two variable columns or a grouping column to proceed.")

def interpret_pval(p):
    if p < 0.01:
        st.success("Strong evidence to reject the null hypothesis (p < 0.01)")
    elif p < 0.05:
        st.info("Moderate evidence to reject the null hypothesis (p < 0.05)")
    elif p < 0.1:
        st.warning("Weak evidence to reject the null hypothesis (p < 0.1)")
    else:
        st.error("No evidence to reject the null hypothesis (p >= 0.1)")

# General Linear Model Tool
def general_linear_model_tool(df):
    st.subheader("General Linear Model (GLM)")

    numeric_columns = df.select_dtypes(include=['number']).columns.tolist()
    categorical_columns = df.select_dtypes(exclude=['number']).columns.tolist()

    y_col = st.selectbox("Select Dependent Variable (Y)", numeric_columns)
    predictors = st.multiselect("Select Independent Variables (X) — can include numeric or categorical", numeric_columns + categorical_columns)

    if st.button("Run General Linear Model"):
        if not y_col or not predictors:
            st.error("Please select both a dependent variable and at least one predictor.")
            return

        formula = f"{y_col} ~ " + " + ".join(predictors)
        st.write(f"Model Formula: `{formula}`")

        try:
            model = smf.ols(formula, data=df).fit()
            st.write("### Model Summary")
            st.text(model.summary())

            st.write("### AI Insights")
            significant_vars = model.pvalues[model.pvalues < 0.05].index.tolist()
            if len(significant_vars) > 1:
                st.success(f"Significant predictors detected: {', '.join([v for v in significant_vars if v != 'Intercept'])}")
            else:
                st.info("No strong predictors found (p < 0.05). Consider reviewing your variable selection.")

        except Exception as e:
            st.error(f"An error occurred: {e}")

# Function to perform Two-Way ANOVA
def two_way_anova(df, dependent_variable, factor1, factor2):
    st.write(f"### Two-Way ANOVA between {dependent_variable}, {factor1}, and {factor2}")
    
    # Perform Two-Way ANOVA
    model = ols(f'{dependent_variable} ~ C({factor1}) + C({factor2}) + C({factor1}):C({factor2})', data=df).fit()
    anova_results = anova_lm(model, typ=2)
    
    st.write(anova_results)
    
    # Interpretation based on p-values
    if anova_results['PR(>F)'][0] < 0.05:
        st.write(f"The **{factor1}** factor has a **statistically significant** effect (p-value < 0.05).")
    if anova_results['PR(>F)'][1] < 0.05:
        st.write(f"The **{factor2}** factor has a **statistically significant** effect (p-value < 0.05).")
    if anova_results['PR(>F)'][2] < 0.05:
        st.write(f"There is a **statistically significant interaction** between {factor1} and {factor2} (p-value < 0.05).")
    
    return anova_results

def interval_plot(df, x_column, y_column):
    st.write(f"### Interval Plot between {x_column} and {y_column}")

    # Plotting interval plot
    fig, ax = plt.subplots(figsize=(10, 6))
    sns.pointplot(x=x_column, y=y_column, data=df, ci="sd", ax=ax)  # Confidence intervals based on standard deviation
    ax.set_title(f"Interval Plot: {y_column} by {x_column}")
    ax.set_xlabel(x_column)
    ax.set_ylabel(y_column)
    st.pyplot(fig)

    # AI Insights
    st.write(f"### AI Insights from the Interval Plot")
    st.write("The plot shows the average values of the dependent variable along with the variability (confidence interval). "
             "If the interval is large, it suggests high variability, while smaller intervals suggest more consistent data.")


# Main Effects Plot function
def main_effects_plot(df, dependent_variable, factor):
    st.write(f"### Main Effects Plot for {dependent_variable} and {factor}")

    # Plotting main effects plot
    fig, ax = plt.subplots(figsize=(10, 6))
    sns.boxplot(x=factor, y=dependent_variable, data=df, ax=ax)
    ax.set_title(f"Main Effects Plot: {dependent_variable} by {factor}")
    ax.set_xlabel(factor)
    ax.set_ylabel(dependent_variable)
    st.pyplot(fig)

    # AI Insights
    st.write(f"### AI Insights from the Main Effects Plot")
    st.write(f"The main effects plot shows how the dependent variable changes with respect to different levels of the factor {factor}. "
             "If there is significant variation between the boxes, it indicates that the factor has a strong impact on the dependent variable.")


# Interactions Plot function
def interactions_plot(df, dependent_variable, factor1, factor2):
    st.write(f"### Interaction Plot between {factor1} and {factor2} on {dependent_variable}")

    # Plotting interactions plot
    fig, ax = plt.subplots(figsize=(10, 6))
    sns.pointplot(x=factor1, y=dependent_variable, hue=factor2, data=df, dodge=True, ci=None, ax=ax)
    ax.set_title(f"Interaction Plot: {factor1} * {factor2} on {dependent_variable}")
    ax.set_xlabel(factor1)
    ax.set_ylabel(dependent_variable)
    st.pyplot(fig)

    # AI Insights
    st.write(f"### AI Insights from the Interaction Plot")
    st.write(f"The interaction plot shows how the relationship between {factor1} and the dependent variable is influenced by the levels of {factor2}. "
             "If the lines are parallel, there is no interaction between the factors. Non-parallel lines indicate a potential interaction.")


# Function to provide insights on the plots
def plot_analysis(df, dependent_variable):
    st.write("### Plot Analysis and Data Insights")

    # Ask for the factors
    categorical_columns = df.select_dtypes(include=['object']).columns.tolist()
    factors = st.multiselect("Select one or two categorical factors for plot analysis", categorical_columns)

    if len(factors) == 1:
        # Interval plot
        interval_plot(df, factors[0], dependent_variable)
    elif len(factors) == 2:
        # Main effects plot and interaction plot
        main_effects_plot(df, dependent_variable, factors[0])
        interactions_plot(df, dependent_variable, factors[0], factors[1])
    else:
        st.warning("Please select at least one categorical variable.")

# Function to perform Repeated Measures ANOVA
def repeated_measures_anova(df, dependent_variable, subject_id, factor):
    st.write(f"### Repeated Measures ANOVA for {dependent_variable} grouped by {factor} and {subject_id}")
    
    # We assume the subject ID represents repeated measures (e.g., same subjects measured multiple times)
    model = ols(f'{dependent_variable} ~ C({factor}) + C({subject_id})', data=df).fit()
    anova_results = anova_lm(model, typ=2)
    
    st.write(anova_results)
    
    # Interpretation
    if anova_results['PR(>F)'][0] < 0.05:
        st.write(f"The **{factor}** factor has a **statistically significant** effect (p-value < 0.05).")
    if anova_results['PR(>F)'][1] < 0.05:
        st.write(f"The **{subject_id}** factor has a **statistically significant** effect (p-value < 0.05).")
    
    return anova_results


# Function to choose the right ANOVA
def perform_anova(df):
    st.write("### ANOVA Analysis")

    # Select dependent variable (continuous data)
    dependent_variable = st.selectbox("Select the dependent variable (continuous data)", df.select_dtypes(include=[np.number]).columns.tolist())
    
    # Select independent variables (categorical data)
    categorical_columns = df.select_dtypes(include=['object']).columns.tolist()
    factors = st.multiselect("Select one or two categorical factors", categorical_columns)
    
    if len(factors) == 1:
        # One-Way ANOVA
        hypothesis_testing_tool(df, dependent_variable, factors[0])
    elif len(factors) == 2:
        # Two-Way ANOVA
        two_way_anova(df, dependent_variable, factors[0], factors[1])
    else:
        st.warning("Please select at least one categorical variable.")
        
    # Repeated Measures ANOVA (if subject is available)
    if "Subject" in df.columns:
        if st.checkbox("Run Repeated Measures ANOVA"):
            repeated_measures_anova(df, dependent_variable, "Subject", factors[0] if len(factors) > 0 else "Factor")

# Scatter Plot Tool with AI Inference
def scatter_plot_ai(df):
    st.subheader("Scatter Plot with AI Insights")

    numeric_columns = df.select_dtypes(include=['number']).columns.tolist()
    categorical_columns = df.select_dtypes(exclude=['number']).columns.tolist()

    x_col = st.selectbox("Select X-axis Column", numeric_columns)
    y_col = st.selectbox("Select Y-axis Column", numeric_columns)
    category = st.selectbox("Select Category Column (optional)", [None] + categorical_columns)

    if st.button("Generate Scatter Plot"):
        fig, ax = plt.subplots()
        if category:
            sns.scatterplot(data=df, x=x_col, y=y_col, hue=category, ax=ax)
        else:
            sns.scatterplot(data=df, x=x_col, y=y_col, ax=ax)
        st.pyplot(fig)

        # AI insights: correlation, regression, clustering
        correlation = df[[x_col, y_col]].corr().iloc[0,1]
        st.write(f"Correlation between {x_col} and {y_col}: {correlation:.2f}")

        lr = LinearRegression()
        clean_df = df[[x_col, y_col]].dropna()
        lr.fit(clean_df[[x_col]], clean_df[y_col])
        st.write(f"Linear trend: y = {lr.coef_[0]:.2f}x + {lr.intercept_:.2f}")

        kmeans = KMeans(n_clusters=2)
        labels = kmeans.fit_predict(clean_df)
        clean_df['Cluster'] = labels
        fig2, ax2 = plt.subplots()
        sns.scatterplot(data=clean_df, x=x_col, y=y_col, hue='Cluster', palette='Set2', ax=ax2)
        st.pyplot(fig2)
        st.write("AI Insight: Clusters identified suggest subgroups or patterns in the data.")

# Individual Value Plot with AI Inference
def individual_value_plot_ai(df):
    st.subheader("Individual Value Plot with AI Insights")

    numeric_columns = df.select_dtypes(include=['number']).columns.tolist()
    categorical_columns = df.select_dtypes(exclude=['number']).columns.tolist()

    y_col = st.selectbox("Select Variable Column (Y-axis)", numeric_columns)
    x_col = st.selectbox("Select Category Column (X-axis)", categorical_columns)

    if st.button("Generate Individual Value Plot"):
        fig, ax = plt.subplots()
        sns.stripplot(data=df, x=x_col, y=y_col, jitter=True, ax=ax)
        st.pyplot(fig)

        group_means = df.groupby(x_col)[y_col].mean()
        st.write("Group Means:")
        st.write(group_means)

        st.write("AI Insight: Differences in mean values across categories suggest possible effects or trends.")

# DOE Setup and Guidance
def design_of_experiment():
    st.write("### Design of Experiment (DOE) Setup")
    
    # Step 1: Ask the user for the factors involved in the experiment
    factors = st.text_area("Enter the factors involved in the experiment (comma separated):")
    factors = [f.strip() for f in factors.split(",")] if factors else []

    # Step 2: Ask the user for the levels of each factor
    levels = st.text_area("Enter the levels for each factor (comma separated):")
    levels = [l.strip() for l in levels.split(",")] if levels else []

    # Step 3: Ask for the response variable
    response_variable = st.text_input("Enter the response variable (the variable to be measured):")

    # Step 4: Generate Experimental Design (Full Factorial or Fractional Factorial)
    if len(factors) > 0 and len(levels) > 0 and response_variable:
        # Generate the factorial design
        design = generate_factorial_design(factors, levels)
        st.write("### Experimental Design:")
        st.write(design)

        # Ask if the user wants to proceed with the experiment setup
        if st.button("Proceed with Experiment Setup"):
            st.write("### Data Collection Instructions")
            st.write(f"To collect the data, you will need to set the following factors at the specified levels: {factors}")
            st.write(f"The response variable you need to measure is: {response_variable}")
            
            # AI guidance for the experiment
            st.write("#### AI Insights:")
            st.write(f"Please collect data based on the experimental design and record the corresponding values of {response_variable}.")
    else:
        st.warning("Please enter valid factors, levels, and response variable.")

        
# Generate Factorial Design based on number of factors and levels
def generate_factorial_design(factors, levels):
    # Create a DataFrame for factorial design
    design = pd.DataFrame(np.array(np.meshgrid(*[levels] * len(factors))).T.reshape(-1, len(factors)), columns=factors)
    
    # Creating a full factorial design (Every combination of factors and levels)
    st.write(f"### Full Factorial Design for {len(factors)} factors and {len(levels)} levels.")
    return design


# Data Analysis of DOE results
def analyze_doe_data(df, factors, response_variable):
    st.write("### DOE Data Analysis")

    # Perform Two-Way ANOVA if two factors are present
    if len(factors) == 2:
        st.write("Performing Two-Way ANOVA...")
        model = ols(f'{response_variable} ~ C({factors[0]}) + C({factors[1]}) + C({factors[0]}):C({factors[1]})', data=df).fit()
        anova_results = anova_lm(model, typ=2)
        st.write(anova_results)
        
        # AI Insights for interpretation of ANOVA results
        if anova_results['PR(>F)'][0] < 0.05:
            st.write(f"Factor **{factors[0]}** has a statistically significant effect.")
        if anova_results['PR(>F)'][1] < 0.05:
            st.write(f"Factor **{factors[1]}** has a statistically significant effect.")
        if anova_results['PR(>F)'][2] < 0.05:
            st.write(f"There is a statistically significant interaction between **{factors[0]}** and **{factors[1]}**.")

    elif len(factors) == 1:
        st.write("Performing One-Way ANOVA...")
        hypothesis_testing_tool(df, response_variable, factors[0])
    
    else:
        st.warning("Please ensure that factors are provided for the analysis.")

# Function to load data from different file formats
def load_data(file):
    if file.name.endswith('.csv'):
        df = pd.read_csv(file)
    elif file.name.endswith('.xlsx') or file.name.endswith('.xls'):
        df = pd.read_excel(file)
    else:
        st.error("Unsupported file type!")
        return None
    return df

# Run charts to check for patterns (Clustering, Mixture, Oscillating, Trend)
def run_charts(df, relevant_column):
    st.write("### Run Charts for Pattern Detection")

    # Time series plot (Run chart)
    st.write("### Run Chart (Raw Data)")
    fig, ax = plt.subplots(figsize=(10, 6))
    ax.plot(df[relevant_column], label='Raw Data')
    ax.set_title('Run Chart (Raw Data)')
    ax.set_xlabel('Index')
    ax.set_ylabel('Value')
    ax.legend()
    st.pyplot(fig)

    # Clustering pattern detection using KMeans
    st.write("### Clustering Pattern Detection")
    kmeans = KMeans(n_clusters=2, random_state=42)
    clusters = kmeans.fit_predict(df[[relevant_column]])
    
    # Plot clusters
    fig, ax = plt.subplots(figsize=(10, 6))
    ax.scatter(df.index, df[relevant_column], c=clusters, cmap='viridis')
    ax.set_title('Clustering Pattern')
    ax.set_xlabel('Index')
    ax.set_ylabel('Value')
    st.pyplot(fig)

    # Mixture pattern detection (e.g., using Gaussian Mixture Model - GMM)
    st.write("### Mixture Pattern Detection")
    from sklearn.mixture import GaussianMixture
    gmm = GaussianMixture(n_components=2, random_state=42)
    gmm_clusters = gmm.fit_predict(df[[relevant_column]])

    # Plot GMM clusters
    fig, ax = plt.subplots(figsize=(10, 6))
    ax.scatter(df.index, df[relevant_column], c=gmm_clusters, cmap='plasma')
    ax.set_title('Mixture Pattern Detection')
    ax.set_xlabel('Index')
    ax.set_ylabel('Value')
    st.pyplot(fig)

    # Trend pattern detection (using ADF test for stationarity)
    st.write("### Trend Pattern Detection (ADF Test for Stationarity)")

    adf_stat, p_value, _, _, _, _ = adfuller(df[relevant_column])
    st.write(f"ADF Statistic: {adf_stat:.4f}")
    st.write(f"P-value: {p_value:.4f}")
    
    if p_value < 0.05:
        st.write("The data shows a **trend** (p-value < 0.05).")
    else:
        st.write("The data **does not show a trend** (p-value > 0.05).")

    # Oscillating pattern detection (look for periodic behavior)
    st.write("### Oscillating Pattern Detection (Periodicity Check)")
    periodogram = np.abs(np.fft.fft(df[relevant_column]))
    
    # Plot periodogram (Frequency domain representation)
    fig, ax = plt.subplots(figsize=(10, 6))
    ax.plot(periodogram)
    ax.set_title('Oscillating Pattern Detection (Periodogram)')
    ax.set_xlabel('Frequency')
    ax.set_ylabel('Amplitude')
    st.pyplot(fig)

# Function to run descriptive statistics
def descriptive_statistics(df):
    st.write("### Descriptive Statistics")
    
    # Check if the column "Measurements" exists
    if 'Measurements' in df.columns:
        data = df['Measurements']
        
        # Calculate basic statistics
        mean = data.mean()
        median = data.median()
        std_dev = data.std()
        min_val = data.min()
        max_val = data.max()
        q1 = data.quantile(0.25)
        q3 = data.quantile(0.75)
        
        # Display results
        st.write(f"Mean: {mean:.4f}")
        st.write(f"Median: {median:.4f}")
        st.write(f"Standard Deviation: {std_dev:.4f}")
        st.write(f"Min: {min_val:.4f}")
        st.write(f"Max: {max_val:.4f}")
        st.write(f"25th Percentile (Q1): {q1:.4f}")
        st.write(f"75th Percentile (Q3): {q3:.4f}")
        
        # Plot the histogram using Matplotlib
        st.write("### Histogram of Measurements")
        plt.figure(figsize=(10, 6))
        plt.hist(data, bins=30, color='lightblue', edgecolor='black')
        plt.title("Histogram of Measurements")
        plt.xlabel('Measurement')
        plt.ylabel('Frequency')
        st.pyplot(plt)  # Display the plot
    else:
        st.error("The dataset does not contain a 'Measurements' column.")


# Six-Piece Capability function (Xbar and R charts, last 25 subgroups, etc.)
def six_piece_capability(df, relevant_column, subgroup_size=5):
    st.write("### Six-Piece Capability Study")
    
    # Create subgroups by grouping the data
    subgroups = [df[relevant_column].iloc[i:i + subgroup_size] for i in range(0, len(df), subgroup_size)]
    subgroups = [s for s in subgroups if len(s) == subgroup_size]  # Only keep complete subgroups
    
    # Calculate X-bar and R
    x_bar = [np.mean(s) for s in subgroups]  # X-bar: mean of each subgroup
    r = [np.max(s) - np.min(s) for s in subgroups]  # R: range of each subgroup
    
    # Plot the X-bar chart
    fig, ax = plt.subplots(figsize=(10, 6))
    ax.plot(x_bar, marker='o', label='X-bar')
    ax.axhline(np.mean(x_bar), color='red', linestyle='dashed', label='Average X-bar')
    ax.set_title('X-bar Chart')
    ax.set_xlabel('Subgroup')
    ax.set_ylabel('X-bar Value')
    ax.legend()
    st.pyplot(fig)

    # Plot the R-chart
    fig, ax = plt.subplots(figsize=(10, 6))
    ax.plot(r, marker='o', label='R')
    ax.axhline(np.mean(r), color='red', linestyle='dashed', label='Average R')
    ax.set_title('R-chart')
    ax.set_xlabel('Subgroup')
    ax.set_ylabel('R Value')
    ax.legend()
    st.pyplot(fig)

    # The last 25 subgroups plot
    st.write("### Last 25 Subgroups")
    last_25_subgroups = subgroups[-25:]  # Get the last 25 subgroups
    fig, ax = plt.subplots(figsize=(10, 6))
    ax.boxplot(last_25_subgroups, patch_artist=True)
    ax.set_title('Last 25 Subgroups')
    ax.set_xlabel('Subgroup')
    ax.set_ylabel('Values')
    st.pyplot(fig)

    # Normal Probability Plot for the relevant column
    st.write("### Normal Probability Plot")
    fig, ax = plt.subplots(figsize=(10, 6))
    stats.probplot(df[relevant_column], dist="norm", plot=ax)
    st.pyplot(fig)

    # Histogram with specification limits
    st.write("### Histogram")
    fig, ax = plt.subplots(figsize=(10, 6))
    df[relevant_column].hist(ax=ax, bins=20, color='lightblue', edgecolor='black')
    ax.set_title(f"Histogram for {relevant_column}")
    ax.set_xlabel('Data')
    ax.set_ylabel('Frequency')
    st.pyplot(fig)

# Visualizations function
def plot_visualizations(df):
    st.write("### Visualizations")
    fig, ax = plt.subplots(figsize=(10, 8))
    df.hist(ax=ax, bins=20)
    st.pyplot(fig)

    # Boxplot
    fig, ax = plt.subplots(figsize=(10, 6))
    sns.boxplot(data=df, ax=ax)
    st.pyplot(fig)

    # Pairplot (only for numerical columns)
    if len(df.select_dtypes(include=[np.number]).columns) > 1:
        pairplot_fig = sns.pairplot(df)
        st.pyplot(pairplot_fig)


# Anomaly detection function
def detect_anomalies(df):
    st.write("### Anomaly Detection")
    
    # Ensure that only numerical columns are used for anomaly detection
    numerical_df = df.select_dtypes(include=[np.number])
    
    scaler = StandardScaler()
    df_scaled = scaler.fit_transform(numerical_df)
    
    # Create the model and fit it
    model = IsolationForest(contamination=0.05)
    anomalies = model.fit_predict(df_scaled)
    
    # Add anomaly results to the DataFrame
    df['anomaly'] = anomalies
    st.write(f"Anomalies detected: {sum(anomalies == -1)} out of {len(df)}")
    
    # Create a scatter plot for anomalies
    fig, ax = plt.subplots(figsize=(10, 6))
    sns.scatterplot(x=df.columns[0], y=df.columns[1], hue='anomaly', data=df, ax=ax, palette='coolwarm')
    ax.set_title('Anomaly Detection')
    
    # Show the plot in Streamlit
    st.pyplot(fig)


# Function to perform Regression Analysis
def regression_analysis(df):
    # Ensure the dataframe has more than one numeric column for regression
    numeric_columns = df.select_dtypes(include=[np.number]).columns.tolist()

    if len(numeric_columns) < 2:
        st.write("### Regression Analysis")
        st.write("Regression analysis is not possible with a single-column dataset. Please upload a dataset with at least two numeric columns to perform the regression.")
        return

    # If there are enough numeric columns, proceed with regression
    st.write("### Regression Analysis")

    try:
        # Assuming the first column as the independent variable (X) and second column as dependent variable (y)
        X = df[numeric_columns[0]].values.reshape(-1, 1)  # Independent variable (first column)
        y = df[numeric_columns[1]].values  # Dependent variable (second column)

        # Initialize the Linear Regression model
        model = LinearRegression()

        # Fit the model
        model.fit(X, y)

        # Predict the y values using the trained model
        y_pred = model.predict(X)

        # Plot the regression
        fig, ax = plt.subplots(figsize=(10, 6))
        ax.scatter(X, y, color='blue')
        ax.plot(X, y_pred, color='red')
        ax.set_title('Linear Regression')
        ax.set_xlabel(numeric_columns[0])
        ax.set_ylabel(numeric_columns[1])
        st.pyplot(fig)

        # R^2 value and Mathematical Model
        st.write(f"R^2 value: {model.score(X, y)}")
        st.write(f"Mathematical Model: y = {model.coef_[0]:.2f} * X + {model.intercept_:.2f}")

    except ValueError as e:
        # Handle non-numeric values or any other error
        st.write(f"An error occurred during regression analysis: {e}")
        st.write("This could be due to non-numeric values in the dataset. Please ensure the data is clean and only contains numeric values for regression.")

# AI-driven Suggestions function
def ai_suggestions(df):
    st.write("### AI Suggestions Based on Correlation Matrix")
    correlation_matrix = df.corr()
    st.write(correlation_matrix)
    highly_correlated = correlation_matrix[abs(correlation_matrix) > 0.8].stack().reset_index()
    highly_correlated = highly_correlated[highly_correlated['level_0'] != highly_correlated['level_1']]
    if not highly_correlated.empty:
        st.write("Highly correlated features (absolute correlation > 0.8):")
        st.write(highly_correlated)
        st.write("Suggested next steps:")
        st.write("- You might want to explore regression models for these features.")
        st.write("- Consider reducing multicollinearity by removing one of the correlated features.")
    else:
        st.write("No strong correlations found. You can try clustering algorithms like KMeans.")


# KMeans Clustering function
def kmeans_clustering(df, num_clusters=3):
    st.write(f"### KMeans Clustering with {num_clusters} Clusters")
    
    # Ensure that only numerical columns are used for clustering
    numerical_df = df.select_dtypes(include=[np.number])
    
    # Standardize the data
    scaler = StandardScaler()
    X_scaled = scaler.fit_transform(numerical_df)
    
    # Apply KMeans clustering
    kmeans = KMeans(n_clusters=num_clusters, random_state=42)
    clusters = kmeans.fit_predict(X_scaled)
    
    # Add the cluster labels to the DataFrame
    df['Cluster'] = clusters
    
    # Create a scatter plot for the first two columns, colored by cluster
    fig, ax = plt.subplots(figsize=(10, 6))
    sns.scatterplot(x=df.columns[0], y=df.columns[1], hue='Cluster', palette="Set2", data=df, ax=ax, s=100)
    ax.set_title(f"KMeans Clustering with {num_clusters} Clusters")
    
    # Display the plot
    st.pyplot(fig)


# Capability Study function
def capability_study(df, relevant_column, USL, LSL):
    st.write("### Capability Study")

    # Select the relevant column
    data = df[relevant_column]

    # Calculate process statistics
    mean = data.mean()
    std_dev = data.std()

    # Calculate capability indices
    Cp = (USL - LSL) / (6 * std_dev)
    Cpk = min((USL - mean) / (3 * std_dev), (mean - LSL) / (3 * std_dev))
    Pp = (USL - LSL) / (6 * data.std())
    Ppk = min((USL - mean) / (3 * data.std()), (mean - LSL) / (3 * data.std()))
    
    # Calculate parts per million (PPM)
    PPM = (1 - Cpk) * 1_000_000
    
    # Calculate capability indices for measurement system
    Cm = (USL - LSL) / (6 * std_dev)  # Cm for measurement system
    Cmk = min((USL - mean) / (3 * std_dev), (mean - LSL) / (3 * std_dev))  # Cmk for measurement system

    # Display the results
    st.write(f"Process Mean: {mean:.4f}")
    st.write(f"Process Standard Deviation: {std_dev:.4f}")
    st.write(f"Cp: {Cp:.4f}")
    st.write(f"Cpk: {Cpk:.4f}")
    st.write(f"Pp: {Pp:.4f}")
    st.write(f"Ppk: {Ppk:.4f}")
    st.write(f"PPM (Parts per Million): {PPM:.0f}")
    st.write(f"Cm: {Cm:.4f}")
    st.write(f"Cmk: {Cmk:.4f}")

    # Plot histogram with specification limits
    fig, ax = plt.subplots(figsize=(10, 6))
    data.hist(ax=ax, bins=20, color='lightblue', edgecolor='black')
    ax.axvline(USL, color='red', linestyle='dashed', linewidth=2, label="USL")
    ax.axvline(LSL, color='green', linestyle='dashed', linewidth=2, label="LSL")
    ax.set_title(f"Histogram for {relevant_column}")
    ax.set_xlabel('Data')
    ax.set_ylabel('Frequency')
    ax.legend()
    st.pyplot(fig)


# Normality Check function with multiple distribution fitting
def fit_distributions(df, relevant_columns=None):
    st.write("### Distribution Fitting Analysis")
    
    # If relevant_columns is not specified, use all numerical columns
    if relevant_columns is None:
        relevant_columns = df.select_dtypes(include=[np.number]).columns.tolist()

    # Select numerical columns only and remove rows with NaN values
    numerical_df = df[relevant_columns].select_dtypes(include=[np.number]).dropna()

    # Check if there are any valid numerical columns
    if numerical_df.empty:
        st.write("No valid numerical columns found in the dataset.")
        return

    # Loop through each numerical column that is relevant
    for column in numerical_df.columns:
        st.write(f"Checking distributions for {column}")
        
        # Initialize results dictionary
        results = {}

        # First, check for normality using the Shapiro-Wilk test
        st.write("### Normality Test (Shapiro-Wilk)")
        stat, p_value = stats.shapiro(df[column])
        st.write(f"Shapiro-Wilk Test p-value: {p_value:.4f}")
        
        if p_value > 0.05:
            st.write(f"The data for **{column}** is **normally distributed** (p-value > 0.05).")
            # Show probability plot for the normal distribution
            st.write("### Probability Plot for Normal Distribution")
            fig, ax = plt.subplots(figsize=(10, 6))
            stats.probplot(df[column], dist="norm", plot=ax)
            st.pyplot(fig)
            continue  # Skip the remaining fitting if it's already normal

        else:
            st.write(f"The data for **{column}** is **not normally distributed** (p-value < 0.05). Proceeding to fit alternative distributions.")

        # Now proceed with fitting other distributions
        try:
            shape, loc, scale = stats.lognorm.fit(df[column], floc=0)
            _, ad_value, p_value = stats.anderson(df[column], dist='norm')
            results['Lognormal'] = {'AD Value': ad_value, 'P-value': p_value}
        except Exception as e:
            st.write(f"Error fitting Lognormal: {e}")
        
        try:
            loc, scale = stats.expon.fit(df[column], floc=0)
            _, ad_value, p_value = stats.anderson(df[column], dist='expon')
            results['Exponential'] = {'AD Value': ad_value, 'P-value': p_value}
        except Exception as e:
            st.write(f"Error fitting Exponential: {e}")

        try:
            shape, loc, scale = stats.weibull_min.fit(df[column], floc=0)
            _, ad_value, p_value = stats.anderson(df[column], dist='weibull_min')
            results['Weibull'] = {'AD Value': ad_value, 'P-value': p_value}
        except Exception as e:
            st.write(f"Error fitting Weibull: {e}")

        try:
            df_boxcox, _ = stats.boxcox(df[column] + 1)  # Adding 1 to avoid log(0)
            _, ad_value, p_value = stats.anderson(df_boxcox, dist='norm')
            results['Box-Cox'] = {'AD Value': ad_value, 'P-value': p_value}
        except Exception as e:
            st.write(f"Error fitting Box-Cox: {e}")

        try:
            shape, loc, scale = stats.lognorm.fit(df[column])
            _, ad_value, p_value = stats.anderson(df[column], dist='norm')
            results['3-Parameter Lognormal'] = {'AD Value': ad_value, 'P-value': p_value}
        except Exception as e:
            st.write(f"Error fitting 3-Parameter Lognormal: {e}")

        # Ensure results are non-empty and contain scalar values
        if results:
            # Filter out non-scalar AD values from results
            valid_results = {k: v for k, v in results.items() if isinstance(v['AD Value'], (int, float))}

            if valid_results:
                # Select the best distribution (lowest AD value)
                best_distribution = min(valid_results, key=lambda x: valid_results[x]['AD Value'])
                best_ad_value = valid_results[best_distribution]['AD Value']
                best_p_value = valid_results[best_distribution]['P-value']

                # Show the results for each distribution
                st.write("Results for each distribution:")
                for dist, result in valid_results.items():
                    st.write(f"{dist} - AD Value: {result['AD Value']:.4f}, P-value: {result['P-value']:.4f}")
                
                # Display the best distribution
                st.write(f"\nBest Distribution: **{best_distribution}**")
                st.write(f"AD Value: {best_ad_value:.4f}, P-value: {best_p_value:.4f}")

                # Explain the choice of the best distribution
                st.write(f"\n### AI Explanation")
                st.write(f"The **{best_distribution}** was chosen because it has the lowest Anderson-Darling (AD) value, indicating the best fit to the data. The AI selected this distribution based on the goodness-of-fit measure provided by the AD statistic.")
                
                # Show Probability Plot for the best distribution
                st.write(f"### Probability Plot for {best_distribution}")
                fig, ax = plt.subplots(figsize=(10, 6))
                if best_distribution == "Lognormal":
                    stats.probplot(df[column], dist="lognorm", plot=ax)
                elif best_distribution == "Exponential":
                    stats.probplot(df[column], dist="expon", plot=ax)
                elif best_distribution == "Weibull":
                    stats.probplot(df[column], dist="weibull_min", plot=ax)
                elif best_distribution == "Box-Cox":
                    stats.probplot(df_boxcox, dist="norm", plot=ax)
                elif best_distribution == "3-Parameter Lognormal":
                    stats.probplot(df[column], dist="lognorm", plot=ax)
                st.pyplot(fig)

            else:
                st.write("No valid distributions could be fitted to the data.")
        else:
            st.write("No distributions were fit successfully.")

        # Control Chart function
def control_chart(df, relevant_column):
    st.write("### Control Chart Analysis")

    # Step 1: Analyze the data type (attribute vs. variable data)
    data_type = st.selectbox("Select Data Type for Control Chart", ["Variable Data", "Attribute Data"])

    if data_type == "Variable Data":
        st.write("For variable data, the AI recommends an **X-bar and R Chart** or **P Chart** depending on the situation.")

        # Step 2: Ask for the number of subgroups (for X-bar and R Chart)
        num_subgroups = st.slider("Select number of subgroups", 2, 10, 5)
        
        # If enough subgroups are available, plot X-bar and R Chart
        if num_subgroups >= 2:
            st.write("### X-bar and R Chart")
            subgroup_data = [df[relevant_column].iloc[i:i + num_subgroups] for i in range(0, len(df), num_subgroups)]
            subgroup_data = [s for s in subgroup_data if len(s) == num_subgroups]  # Only keep complete subgroups

            x_bar = [np.mean(s) for s in subgroup_data]
            r = [np.max(s) - np.min(s) for s in subgroup_data]
            
            # Plot X-bar Chart
            fig, ax = plt.subplots(figsize=(10, 6))
            ax.plot(x_bar, marker='o', label='X-bar')
            ax.axhline(np.mean(x_bar), color='red', linestyle='dashed', label='Average X-bar')
            ax.set_title('X-bar Chart')
            ax.set_xlabel('Subgroup')
            ax.set_ylabel('X-bar Value')
            ax.legend()
            st.pyplot(fig)

            # Plot R Chart
            fig, ax = plt.subplots(figsize=(10, 6))
            ax.plot(r, marker='o', label='R')
            ax.axhline(np.mean(r), color='red', linestyle='dashed', label='Average R')
            ax.set_title('R Chart')
            ax.set_xlabel('Subgroup')
            ax.set_ylabel('R Value')
            ax.legend()
            st.pyplot(fig)

            # AI Insights
            st.write(f"### AI Insights from the Control Chart")
            st.write("The X-bar and R charts are used to monitor the mean and the range of subgroups. If the points fall outside the control limits, it may indicate an issue with the process stability.")
        
    elif data_type == "Attribute Data":
        st.write("For attribute data, the AI will suggest a **P Chart**, **NP Chart**, **C Chart**, or **U Chart** depending on your data.")

        # Ask for specific chart selection
        chart_type = st.selectbox("Select the type of Attribute Data Control Chart", ["P Chart", "NP Chart", "C Chart", "U Chart"])
        
        if chart_type == "P Chart":
            st.write("### P Chart for Proportion Defective")
            # Assume df[relevant_column] contains the proportions of defective items
            defect_counts = df[relevant_column].value_counts(normalize=True)
            fig, ax = plt.subplots(figsize=(10, 6))
            ax.plot(defect_counts.index, defect_counts.values, marker='o', label='Defective Proportion')
            ax.set_title('P Chart')
            ax.set_xlabel('Sample')
            ax.set_ylabel('Proportion Defective')
            ax.legend()
            st.pyplot(fig)

            st.write(f"### AI Insights for the P Chart")
            st.write("The P Chart helps monitor the proportion of defective items. A sudden increase in defect proportion may indicate a problem with the process.")

        elif chart_type == "NP Chart":
            st.write("### NP Chart for Number of Defective Items")
            # Assume df[relevant_column] contains the number of defectives in each sample
            defect_counts = df[relevant_column].value_counts()
            fig, ax = plt.subplots(figsize=(10, 6))
            ax.plot(defect_counts.index, defect_counts.values, marker='o', label='Defective Count')
            ax.set_title('NP Chart')
            ax.set_xlabel('Sample')
            ax.set_ylabel('Number of Defectives')
            ax.legend()
            st.pyplot(fig)

            st.write(f"### AI Insights for the NP Chart")
            st.write("The NP Chart monitors the count of defective items in a sample. If defect counts increase, it may signal an issue with production or quality control.")

        elif chart_type == "C Chart":
            st.write("### C Chart for Number of Defects")
            # Assume df[relevant_column] contains the number of defects in each sample
            defect_counts = df[relevant_column].value_counts()
            fig, ax = plt.subplots(figsize=(10, 6))
            ax.plot(defect_counts.index, defect_counts.values, marker='o', label='Defects')
            ax.set_title('C Chart')
            ax.set_xlabel('Sample')
            ax.set_ylabel('Defects')
            ax.legend()
            st.pyplot(fig)

            st.write(f"### AI Insights for the C Chart")
            st.write("The C Chart monitors the number of defects per unit in a sample. A sudden increase in the defect count may suggest a quality issue.")

        elif chart_type == "U Chart":
            st.write("### U Chart for Defects per Unit")
            # Assume df[relevant_column] contains defects per unit
            defects_per_unit = df[relevant_column]
            fig, ax = plt.subplots(figsize=(10, 6))
            ax.plot(defects_per_unit, marker='o', label='Defects per Unit')
            ax.set_title('U Chart')
            ax.set_xlabel('Sample')
            ax.set_ylabel('Defects per Unit')
            ax.legend()
            st.pyplot(fig)

            st.write(f"### AI Insights for the U Chart")
            st.write("The U Chart monitors the number of defects per unit. If the number of defects per unit increases, this may indicate an issue with the production process.")

# Function to perform Multi-Vari Analysis (Plot all individual data points)
def multi_vari_analysis(df):
    st.write("### Multi-Vari Analysis")

    # Identify the continuous and categorical variables
    numeric_columns = df.select_dtypes(include=[np.number]).columns.tolist()
    categorical_columns = df.select_dtypes(include=['object']).columns.tolist()

    # Ask user to select a continuous variable for analysis
    selected_continuous = st.selectbox("Select continuous variable for multi-vari analysis", numeric_columns)

    # If there are categorical variables, ask the user to select one
    selected_categorical = None
    if len(categorical_columns) > 0:
        selected_categorical = st.selectbox("Select categorical variable for multi-vari analysis", categorical_columns)

    # Generate the graph based on the user's selections
    if selected_categorical:
        st.write(f"### Line Chart for {selected_continuous} by {selected_categorical}")

        # Plot each individual data point against the categorical variable
        plt.figure(figsize=(10, 6))
        sns.lineplot(x=selected_categorical, y=selected_continuous, data=df, marker='o', ci=None, markersize=8, linestyle='-', hue=selected_categorical)
        plt.title(f"Line Chart: {selected_continuous} by {selected_categorical}")
        plt.xlabel(selected_categorical)
        plt.ylabel(selected_continuous)
        st.pyplot(plt)

        # AI Insights
        st.write(f"### AI Insights:")
        st.write(f"The line chart shows each individual data point of {selected_continuous} against {selected_categorical}. "
                 "Look for trends, outliers, or clusters of data points in each category.")
    
    else:
        st.write(f"### Simple Line Chart for {selected_continuous}")

        # Plot a simple line chart if no categorical variable is selected
        plt.figure(figsize=(10, 6))
        sns.lineplot(x=df.index, y=selected_continuous, data=df, marker='o', ci=None, markersize=8, linestyle='-')
        plt.title(f"Line Chart: {selected_continuous}")
        plt.xlabel("Index")
        plt.ylabel(selected_continuous)
        st.pyplot(plt)

        # AI Insights
        st.write(f"### AI Insights:")
        st.write(f"The simple line chart shows how {selected_continuous} changes across the data. Look for trends, peaks, and drops in the data.")

# Feature 1: CG and CGK Calculation (no need for Part, Operator, Measurement columns)
def calculate_cg_cgk(df):
    st.write("### CG and CGK Calculation")

    # Ask the user to select the column containing the data
    numeric_columns = df.select_dtypes(include=[np.number]).columns.tolist()
    selected_continuous = st.selectbox("Select continuous variable for CG and CGK calculation", numeric_columns)

    # Get specification limits from the user
    usl = st.number_input("Enter the Upper Specification Limit (USL)", value=10.0)
    lsl = st.number_input("Enter the Lower Specification Limit (LSL)", value=0.0)

    # Calculate the standard deviation of the process data
    process_std = df[selected_continuous].std()

    # Calculate CG and CGK
    cg = (usl - lsl) / (6 * process_std)
    cgk = min((usl - df[selected_continuous].mean()) / (3 * process_std), (df[selected_continuous].mean() - lsl) / (3 * process_std))

    # Display results
    st.write(f"### CG (Capability Index): {cg:.4f}")
    st.write(f"### CGK (Capability Index for Measurement System): {cgk:.4f}")

    # Visualizing the results
    fig, ax = plt.subplots(figsize=(8, 6))
    ax.bar(['CG', 'CGK'], [cg, cgk])
    ax.set_title('CG and CGK Values')
    st.pyplot(fig)

# Feature 2: GRR (Gauge Repeatability and Reproducibility)
def grr_calculation(df):
    st.write("### GRR (Gauge Repeatability and Reproducibility) Calculation")

    # Ensure the dataset has 'Part', 'Operator', and 'Measurement' columns
    if 'Part' not in df.columns or 'Operator' not in df.columns or 'Measurement' not in df.columns:
        st.error("The dataset must contain 'Part', 'Operator', and 'Measurement' columns.")
        return
    
    # Grouping data by parts and operators
    grouped_data = df.groupby(['Part', 'Operator']).agg({'Measurement': ['mean', 'std', 'count']}).reset_index()

    # Calculate GRR components (repeatability and reproducibility)
    # MS within (Repeatability)
    ms_within = grouped_data['Measurement']['std'].var()

    # MS between (Reproducibility)
    ms_between = grouped_data.groupby('Part')['Measurement']['mean'].std()

    # Repeatability (Equipment Variation) - within
    repeatability = np.sqrt(ms_within)

    # Reproducibility (Appraiser Variation) - between
    reproducibility = np.sqrt(ms_between)

    # Total GRR (Repeatability + Reproducibility)
    total_grr = repeatability + reproducibility

    # Display results
    st.write(f"### Repeatability (Equipment Variation): {repeatability:.4f}")
    st.write(f"### Reproducibility (Appraiser Variation): {reproducibility:.4f}")
    st.write(f"### Total GRR: {total_grr:.4f}")

    # Visualizing GRR results using a bar chart
    fig, ax = plt.subplots(figsize=(8, 6))
    ax.bar(['Repeatability', 'Reproducibility', 'Total GRR'], [repeatability, reproducibility, total_grr])
    ax.set_title('GRR Breakdown')
    st.pyplot(fig)

# Time Series Analysis Feature
def time_series_analysis(df):
    st.write("### Time Series Analysis")

    # Identify the continuous variables (potential time series columns)
    numeric_columns = df.select_dtypes(include=[np.number]).columns.tolist()
    selected_column = st.selectbox("Select time series variable", numeric_columns)

    # Ask user for the analysis type
    analysis_type = st.selectbox("Select type of analysis", ["Trend", "Seasonality", "Forecasting", "Anomaly Detection"])

    # Time series data (must be a column with continuous data)
    time_series_data = df[selected_column]

    # Perform Time Series Analysis based on the selected option
    if analysis_type == "Trend":
        st.write(f"### Trend Analysis for {selected_column}")
        plt.figure(figsize=(10, 6))
        plt.plot(time_series_data)
        plt.title(f"Trend Analysis of {selected_column}")
        plt.xlabel("Time")
        plt.ylabel(selected_column)
        st.pyplot(plt)

    elif analysis_type == "Seasonality":
        st.write(f"### Seasonality Analysis for {selected_column}")
        plt.figure(figsize=(10, 6))
        sns.lineplot(x=df.index, y=selected_column, data=df)
        plt.title(f"Seasonality in {selected_column}")
        plt.xlabel("Time")
        plt.ylabel(selected_column)
        st.pyplot(plt)

    elif analysis_type == "Forecasting":
        st.write(f"### Forecasting for {selected_column}")
        forecast_periods = st.number_input("Enter the number of periods to forecast", min_value=1, value=12)

        # Using Exponential Smoothing for forecasting
        model = ExponentialSmoothing(time_series_data, trend='add', seasonal='add', seasonal_periods=12)
        model_fit = model.fit()
        forecast = model_fit.forecast(forecast_periods)

        # Plotting the forecasted values
        plt.figure(figsize=(10, 6))
        plt.plot(time_series_data, label='Original Data')
        plt.plot(range(len(time_series_data), len(time_series_data) + forecast_periods), forecast, label='Forecast', linestyle='--')
        plt.title(f"Forecasting {selected_column}")
        plt.xlabel("Time")
        plt.ylabel(selected_column)
        plt.legend()
        st.pyplot(plt)

        # Error Metrics
        st.write(f"### Forecasting Error Metrics")
        mape = mean_absolute_percentage_error(time_series_data[-forecast_periods:], forecast)
        mad = mean_absolute_error(time_series_data[-forecast_periods:], forecast)
        msd = mean_squared_error(time_series_data[-forecast_periods:], forecast)

        st.write(f"MAPE: {mape:.4f}")
        st.write(f"MAD: {mad:.4f}")
        st.write(f"MSD: {msd:.4f}")

    elif analysis_type == "Anomaly Detection":
        st.write(f"### Anomaly Detection for {selected_column}")

        # Detect anomalies using Z-Score (or another method, like Isolation Forest, depending on requirements)
        z_scores = (time_series_data - time_series_data.mean()) / time_series_data.std()
        anomalies = time_series_data[abs(z_scores) > 3]  # Z-Score threshold for anomalies

        plt.figure(figsize=(10, 6))
        plt.plot(time_series_data, label="Original Data")
        plt.scatter(anomalies.index, anomalies, color='red', label="Anomalies")
        plt.title(f"Anomaly Detection in {selected_column}")
        plt.xlabel("Time")
        plt.ylabel(selected_column)
        plt.legend()
        st.pyplot(plt)

        st.write(f"### Detected Anomalies:")
        st.write(anomalies)

def interval_plot(df):
    st.write("### Interval Plot (Showing Min and Max Ranges)")

    # Identify numerical columns in the dataframe
    numeric_columns = df.select_dtypes(include=[np.number]).columns.tolist()

    if len(numeric_columns) == 0:
        st.write("No numerical columns found in the dataset.")
        return

    # Ask the user to select the numerical column for analysis
    selected_numeric_column = st.selectbox("Select a numerical column for the interval plot", numeric_columns)

    # Identify categorical columns in the dataframe
    categorical_columns = df.select_dtypes(include=['object', 'category']).columns.tolist()

    if len(categorical_columns) == 0:
        st.write("No categorical columns found in the dataset.")
        return

    # Ask the user to select the categorical column for analysis
    selected_category_column = st.selectbox("Select a categorical column for the interval plot", categorical_columns)

    if selected_category_column and selected_numeric_column:
        # Group the data by the selected categorical column and calculate the min and max for the numeric column
        grouped = df.groupby(selected_category_column)[selected_numeric_column].agg([min, max]).reset_index()
        
        # Create an interval plot
        fig, ax = plt.subplots(figsize=(10, 6))
        ax.plot(grouped[selected_category_column], grouped['min'], marker='o', label="Min", linestyle='-', color='blue')
        ax.plot(grouped[selected_category_column], grouped['max'], marker='o', label="Max", linestyle='-', color='red')
        
        ax.set_title(f"Interval Plot: {selected_numeric_column} by {selected_category_column}")
        ax.set_xlabel(selected_category_column)
        ax.set_ylabel(selected_numeric_column)
        ax.legend()
        
        # Show the plot
        st.pyplot(fig)

        # AI Insights
        st.write(f"### AI Insights from the Interval Plot")
        st.write(f"The plot shows the minimum (blue) and maximum (red) values for each category in {selected_category_column}.")
        st.write("A wider range between the minimum and maximum indicates more variability, while a narrower range suggests more consistency within the category.")
        st.write("Look for any significant outliers or categories with particularly wide or narrow ranges.")

import pandas as pd
import matplotlib.pyplot as plt
import numpy as np
import streamlit as st

# Function to generate Pareto chart
def generate_pareto(df, category_column):
    st.write(f"### Pareto Chart for {category_column}")

    # Count the occurrences of each category
    category_counts = df[category_column].value_counts()

    # Calculate cumulative percentage
    category_counts = category_counts.sort_values(ascending=False)
    cumulative_percentage = category_counts.cumsum() / category_counts.sum() * 100

    # Plot Pareto chart
    fig, ax = plt.subplots(figsize=(10, 6))
    ax.bar(category_counts.index, category_counts.values, color='b', alpha=0.7)
    ax.set_xlabel(category_column)
    ax.set_ylabel('Frequency')
    ax.set_title(f'Pareto Chart for {category_column}')

    ax2 = ax.twinx()
    ax2.plot(category_counts.index, cumulative_percentage.values, color='r', marker='o', linestyle='dashed')
    ax2.set_ylabel('Cumulative Percentage', color='r')
    ax2.tick_params(axis='y', labelcolor='r')

    plt.xticks(rotation=45)
    st.pyplot(fig)

# Function to provide AI-generated insights
def generate_insights(df, category_column):
    st.write(f"### Insights for {category_column}")

    insights = []

    # Basic statistics: mean, median, mode
    category_counts = df[category_column].value_counts()

    most_common = category_counts.idxmax()
    least_common = category_counts.idxmin()
    average_count = category_counts.mean()
    median_count = category_counts.median()

    insights.append(f'Most common category: {most_common}')
    insights.append(f'Least common category: {least_common}')
    insights.append(f'Average count across categories: {average_count:.2f}')
    insights.append(f'Median count across categories: {median_count:.2f}')

    # Checking for potential trends (top 3 categories with largest counts)
    top_categories = category_counts.head(3)
    insights.append('Top 3 categories by count:')
    for category, count in top_categories.items():
        insights.append(f'  {category}: {count}')

    for insight in insights:
        st.write(f"- {insight}")

# Main function for Pareto Analysis in Streamlit format
def pareto_analysis_tool(df):
    st.write("### Pareto Analysis Tool")

    # Identify the categorical columns
    categorical_columns = df.select_dtypes(include=['object']).columns.tolist()

    # User selects the category column for analysis
    selected_column = st.selectbox("Select category column for Pareto analysis", categorical_columns)

    # Generate Pareto chart and insights
    generate_pareto(df, selected_column)
    generate_insights(df, selected_column)

# Function to run Streamlit interface
def run_streamlit_interface():
    st.title("Auto-Mode Statistics Software")

    if st.button("Run Shainin Red X® Methodology"):
        run_redx_tool()

    uploaded_file = st.file_uploader("Choose a file", type=["csv", "xlsx", "xls"])

    if uploaded_file:
        df = load_data(uploaded_file)
        if df is not None:
            st.write("### Data Overview:")
            st.write(df.head())

            # Create clickable menu options
            menu_options = ["Basic Statistics", "Quality", "Plot", "MSA", "ANOVA", "DOE", "Machine Learning", "Time Series"]

            menu_choice = st.selectbox("Select an Analysis Category", menu_options)

            # Display sub-options based on menu choice
            if menu_choice == "Basic Statistics":
                st.write("### Basic Statistics Options")
                if st.button("Run Descriptive Statistics"):
                    descriptive_statistics(df)
                if st.button("Detect Anomalies"):
                    detect_anomalies(df)
                if st.button("Run Regression Analysis"):
                    regression_analysis(df)

            elif menu_choice == "Quality":
                st.write("### Quality Control Options")
                if st.button("Check Normality and Fit Distributions"):
                    fit_distributions(df)
                if st.button("Run Capability Study"):
                    capability_study(df)
                if st.button("Run Six-Piece Capability Study"):
                    six_piece_capability(df)
                if st.button("Run Charts for Clustering, Mixture, Oscillating, and Trend Patterns"):
                    run_charts(df)
                if st.button("Run Control Chart Analysis"):
                    control_chart(df)
                if st.button("Run Multi-Vari Analysis"):
                    multi_vari_analysis(df)
                if st.button("Run Pareto Analysis"):
                    pareto_analysis_tool(df)

            elif menu_choice == "Plot":
                st.write("### Plotting Options")
                if st.button("Generate Visualizations"):
                    plot_visualizations(df)
                if st.button("Run Interval Plot Analysis"):
                    interval_plot(df)
                if st.button("Generate Scatter Plot"):
                    scatter_plot_ai(df)
                if st.button("Generate Individual Value Plot"):
                    individual_value_plot_ai(df)

            elif menu_choice == "MSA":
                st.write("### MSA Options")
                if st.button("Run CG and CGK Calculation"):
                    calculate_cg_cgk(df)
                if st.button("Run GRR Calculation"):
                    grr_calculation(df)

            elif menu_choice == "ANOVA":
                st.write("### ANOVA Options")
                if st.button("Run ANOVA Analysis"):
                    perform_anova(df)
                if st.button("Run Plot Analysis (Interval, Main Effects, and Interactions)"):
                    plot_analysis(df)
                if st.button("Run Hypothesis Testing"):
                    hypothesis_testing_tool(df)
                if st.button("Run General Linear Model"):
                    general_linear_model_tool(df)

            elif menu_choice == "DOE":
                st.write("### DOE Options")
                if st.button("Run Design of Experiment (DOE) Setup"):
                    design_of_experiment()
                if st.button("Analyze DOE Results"):
                    analyze_doe_data(df)

            elif menu_choice == "Machine Learning":
                st.write("### Machine Learning Options")
                num_clusters = st.slider("Select number of clusters", 2, 10, 3)
                if st.button("Run KMeans Clustering"):
                    kmeans_clustering(df, num_clusters)

            elif menu_choice == "Time Series":
                st.write("### Time Series Analysis Options")
                if st.button("Run Time Series Analysis"):
                    time_series_analysis(df)
    
    # After selecting DMAIC insight, proceed with 8D generation
    if st.button("Generate 8D Report"):
        generate_8d_report(df, dmaic_insights=None, selected_analysis=None)

# Running the interface
if __name__ == "__main__":
    run_streamlit_interface()


