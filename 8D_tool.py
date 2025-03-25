import streamlit as st
import random
import matplotlib.pyplot as plt
import io
from PIL import Image
import datetime
import pandas as pd
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches
import os

# ---------------------------
# Simulated database of roles by expertise area
# ---------------------------
ROLE_DATABASE = {
    "quality": [
        {"role": "Quality Engineer", "profile": "Expert in root cause analysis and defect resolution"},
        {"role": "QA Technician", "profile": "Experienced in inspection and interim containment actions"},
    ],
    "production": [
        {"role": "Production Supervisor", "profile": "Manages operators and processes on the shop floor"},
        {"role": "Process Engineer", "profile": "Optimizes manufacturing processes and workflows"},
    ],
    "maintenance": [
        {"role": "Maintenance Technician", "profile": "Handles equipment repairs and failure diagnostics"},
        {"role": "Reliability Engineer", "profile": "Focuses on reducing downtime and chronic failures"},
    ],
    "design": [
        {"role": "Product Design Engineer", "profile": "Specializes in part functionality and failure prevention"},
        {"role": "CAD Specialist", "profile": "Provides technical drawings and design updates"},
    ],
    "safety": [
        {"role": "EHS Coordinator", "profile": "Ensures safety compliance and risk mitigation"},
        {"role": "Safety Officer", "profile": "Handles workplace incident response and audits"},
    ]
}

# ---------------------------
# Keywords to detect department involvement based on problem description
# ---------------------------
KEYWORDS = {
    "defect": "quality",
    "scrap": "quality",
    "inspection": "quality",
    "machine": "maintenance",
    "downtime": "maintenance",
    "line": "production",
    "process": "production",
    "design": "design",
    "blueprint": "design",
    "hazard": "safety",
    "injury": "safety",
    "safety": "safety"
}

# ---------------------------
# Title and Section: D1 – Team Formation
# ---------------------------
st.title("🛠️ 8D Problem Solving Tool")
st.header("D1 – Team Formation")

problem_description = st.text_area("Describe the problem:", help="Be specific about what's happening, where, and with what kind of part or process.")

if st.button("🔍 Get Recommended Roles"):
    found_roles = set()
    description_lower = problem_description.lower()
    for keyword, dept in KEYWORDS.items():
        if keyword in description_lower:
            for role in ROLE_DATABASE[dept]:
                found_roles.add((role["role"], role["profile"]))

    if found_roles:
        st.subheader("👥 Recommended Team Roles:")
        for role, profile in found_roles:
            st.markdown(f"- **{role}** – {profile}")
    else:
        st.info("No specific roles found. Please provide more detail in your problem description.")

st.markdown("---")
st.subheader("👤 Assign Team Members")

num_members = st.number_input("How many team members will be part of this 8D team?", min_value=1, max_value=10, step=1)
team = []
for i in range(num_members):
    st.markdown(f"#### Team Member {i+1}")
    name = st.text_input(f"Name (Member {i+1})", key=f"name_{i}")
    role = st.text_input(f"Role (Member {i+1})", key=f"role_{i}")
    responsibility = st.text_area(f"Responsibility (Member {i+1})", key=f"resp_{i}")
    team.append({"name": name, "role": role, "responsibility": responsibility})

if st.button("✅ Confirm Team"):
    st.success("Team assigned successfully!")
    st.write("### 🧾 Team Overview")
    for member in team:
        st.markdown(f"- **{member['name']}** – *{member['role']}*: {member['responsibility']}")

st.caption("D1: Define and assign a cross-functional team to solve the problem.")

# ---------------------------
# Function: AI suggestion generator for Fishbone categories
# ---------------------------
def generate_ai_fishbone_suggestions(problem_text):
    suggestions = {
        "Man": "Lack of training or unclear instructions.",
        "Machine": "Frequent breakdowns or uncalibrated equipment.",
        "Method": "Non-standardized work procedures.",
        "Material": "Poor material quality or incorrect part used.",
        "Measurement": "Inaccurate measurements or tools.",
        "Environment": "Temperature or lighting affecting results."
    }
    if "operator" in problem_text.lower():
        suggestions["Man"] = "Operator unfamiliar with the updated work instruction."
    if "sensor" in problem_text.lower():
        suggestions["Measurement"] = "Sensor drift causing inaccurate data."
    if "motor" in problem_text.lower():
        suggestions["Machine"] = "Motor overheating during operation."
    return suggestions

# ---------------------------
# Function: AI suggestion generator for 5W2H
# ---------------------------
def generate_ai_5w2h(problem_text):
    return {
        "What": "Defective parts detected during final inspection.",
        "Where": "Assembly line 3 in the main plant.",
        "When": "Observed during the third shift on weekdays.",
        "Who": "Line operators and quality inspectors.",
        "Why": "Increased defect rate affecting customer satisfaction.",
        "How": "Process deviation due to equipment malfunction.",
        "How Many": "Approximately 12% of output from affected line."
    }

# ---------------------------
# Function: AI suggestion generator for Containment Actions
# ---------------------------
def generate_ai_containment_actions(problem_text):
    suggestions = [
        "Isolate all defective parts from good inventory.",
        "Conduct 100% inspection on recent production batches.",
        "Stop production line temporarily to prevent further defects.",
        "Notify quality and production teams for urgent review."
    ]
    if "leak" in problem_text.lower():
        suggestions.append("Seal affected areas and monitor for further leakage.")
    if "overheat" in problem_text.lower():
        suggestions.append("Install temporary cooling or increase monitoring frequency.")
    return suggestions

# ---------------------------
# Section: D2 – Problem Description (AI-enhanced with 5W2H and Fishbone)
# ---------------------------
st.markdown("---")
st.header("D2 – Problem Description")

problem_description_d2 = st.text_area("🛠️ Describe the problem in detail (for AI analysis):", key="problem_d2")

if problem_description_d2:
    ai_5w2h = generate_ai_5w2h(problem_description_d2)
    ai_fishbone = generate_ai_fishbone_suggestions(problem_description_d2)

    st.markdown("### 📋 AI-Suggested Problem Statement (5W2H)")
    what = st.text_input("What is the problem?", value=ai_5w2h["What"])
    where = st.text_input("Where does it occur?", value=ai_5w2h["Where"])
    when = st.text_input("When does it happen?", value=ai_5w2h["When"])
    who = st.text_input("Who is involved?", value=ai_5w2h["Who"])
    why = st.text_input("Why is it a problem?", value=ai_5w2h["Why"])
    how = st.text_input("How is the problem occurring?", value=ai_5w2h["How"])
    how_many = st.text_input("How many parts/products are affected?", value=ai_5w2h["How Many"])

    if st.button("📝 Confirm Problem Statement"):
        st.subheader("✅ Structured Problem Statement")
        st.markdown(f"- **What:** {what}")
        st.markdown(f"- **Where:** {where}")
        st.markdown(f"- **When:** {when}")
        st.markdown(f"- **Who:** {who}")
        st.markdown(f"- **Why:** {why}")
        st.markdown(f"- **How:** {how}")
        st.markdown(f"- **How Many:** {how_many}")

    st.markdown("### 🤖 AI-Suggested Causes for Fishbone Diagram")
    causes = {}
    for category, suggestion in ai_fishbone.items():
        causes[category] = st.text_area(f"{category} (AI suggested)", value=suggestion, key=f"ai_{category}")

    if st.button("📈 Generate Fishbone Diagram"):
        fig, ax = plt.subplots(figsize=(10, 6))
        ax.axis('off')
        ax.set_title('Fishbone Diagram', fontsize=16, weight='bold')
        categories = list(causes.keys())
        y_positions = [5, 4, 3, 2, 1, 0]
        for y, cat in zip(y_positions, categories):
            ax.annotate(cat, xy=(0.1, y), fontsize=12, weight='bold')
            ax.plot([0.2, 0.8], [y, 3], color='gray', linewidth=1)
            ax.annotate(causes[cat], xy=(0.82, y), fontsize=10, va='center')
        st.pyplot(fig)

    # ---------------------------
    # Section: D3 – Interim Containment Actions with AI suggestions
    # ---------------------------
    st.markdown("---")
    st.header("D3 – Interim Containment Actions")

    ai_containment = generate_ai_containment_actions(problem_description_d2)
    containment_actions = []

    st.markdown("### 🔧 AI-Suggested Containment Actions")
    num_suggestions = len(ai_containment)
    for i in range(num_suggestions):
        action = st.text_area(f"Containment Action {i+1}", value=ai_containment[i], key=f"containment_{i}")
        containment_actions.append(action)

    st.markdown("You may modify the suggestions or add more actions below:")
    extra_actions_count = st.number_input("How many additional actions would you like to add?", min_value=0, max_value=10, step=1, key="extra_contain")
    for j in range(extra_actions_count):
        extra = st.text_area(f"Additional Action {j+1}", key=f"extra_contain_{j}")
        containment_actions.append(extra)

    if st.button("✅ Save Containment Actions"):
        st.success("Containment actions saved.")
        st.write("### 🧾 Interim Containment Action Summary")
        for idx, act in enumerate(containment_actions):
            st.markdown(f"**{idx+1}.** {act}")

st.caption("D2 & D3: AI-assisted problem definition, cause analysis, and containment planning.")

# ---------------------------
# Section: D4 – Root Cause Analysis (AI-Enhanced)
# ---------------------------
st.markdown("---")
st.header("D4 – Root Cause Analysis")

if problem_description_d2:
    st.subheader("🧠 AI-Generated Root Cause Insights")

    st.markdown("### 🔧 Background & Context")
    st.markdown("Our Ultimate Problem Solving Tool and DMAIC Tool are designed to structure and optimize problem solving.")
    st.markdown("The issue described may relate to quality deviation, machine reliability, or material variability. According to multiple engineering sources and studies in manufacturing physics, such problems often originate in variation introduced through machine wear, environmental instability, or human error.")

    st.markdown("### 🔍 Possible Root Causes (based on your description)")
    st.markdown("- Inconsistent component performance (e.g., due to thermal drift or mechanical play)")
    st.markdown("- Environmental changes influencing process repeatability")
    st.markdown("- Inadequate process control or calibration drift")
    st.markdown("- Human-induced variance through non-standard procedures")

    st.markdown("### 🧪 What to Measure")
    st.markdown("- Measure temperatures, vibrations, humidity, torque, or any critical machine parameter involved.")
    st.markdown("- Analyze parts before and after failure for deformation, discoloration, or resistance.")
    st.markdown("- Compare Y1 (output variable) to multiple Y2s (candidate causes) using our Statistical Software.")

    st.markdown("### 🧰 Methodologies You Can Use")

    st.markdown("#### ✅ Red X Methodology")
    st.markdown("1. Separate BOB (Best of Best) and WOW (Worst of Worst) parts.")
    st.markdown("2. Define Y₀ and verify the measurement system can discriminate between BOB and WOW.")
    st.markdown("3. Build an Is/Is Not Diagram.")
    st.markdown("4. Perform contrast analysis to isolate sources.")
    st.markdown("5. Swap components, dissect subsystems, and test outcomes.")

    st.markdown("#### ✅ Six Sigma DMAIC")
    st.markdown("- **Define**: Use the D2 structured inputs.")
    st.markdown("- **Measure**: Quantify defect rate, cycle time, etc.")
    st.markdown("- **Analyze**: Use Fishbone + 5 Whys.")
    st.markdown("- **Improve**: Run component swaps, parameter optimizations.")
    st.markdown("- **Control**: Set up statistical monitoring using our tools.")

    st.markdown("#### ✅ Other Root Cause Tools")
    st.markdown("- **Component Swap** – exchange parts between BOB/WOW.")
    st.markdown("- **Process Search** – isolate which step introduces variation.")
    st.markdown("- **Dissection (Component Search)** – break down components layer by layer.")
    st.markdown("- **Y1 vs Y2 Correlation** – chart performance vs machine variables.")
    st.markdown("- **BOB vs WOW Comparison** – contrast every measurable feature.")

    st.markdown("You may now use the tool below once you're ready to apply 5 Whys to drill down the cause.")

# Interactive 5 Whys Input
five_whys = []
num_whys = st.number_input("How many levels of 'Why?' would you like to explore?", min_value=3, max_value=7, step=1, key="num_whys")

for i in range(num_whys):
    response = st.text_input(f"Why {i+1}?", key=f"why_{i}")
    five_whys.append(response)

confirmed_root_cause = ""
if st.button("🔍 Summarize Root Cause"):
    st.subheader("🧩 Root Cause Summary")
    for i, reason in enumerate(five_whys):
        st.markdown(f"**Why {i+1}:** {reason}")

    if five_whys[-1]:
        summary_text = f"Based on the 5 Whys analysis, the root cause appears to be: {five_whys[-1]}. This issue likely originated from cascading factors leading to this specific failure. Further validation through component testing or statistical comparison is advised."
        editable_summary = st.text_area("✏️ AI Summary of Root Cause (editable)", value=summary_text, key="editable_root")
        if st.button("✅ Confirm Root Cause"):
            st.session_state.confirmed_root_cause = editable_summary
            st.success("Root cause confirmed and saved.")

# ---------------------------
# Section: D5 – Permanent Corrective Actions (AI-enhanced)
# ---------------------------
st.markdown("---")
st.header("D5 – Permanent Corrective Actions")

corrective_actions = []
if 'confirmed_root_cause' in st.session_state:
    st.markdown("### 🧠 AI-Suggested Corrective Actions")
    st.markdown(f"Based on the confirmed root cause: *{st.session_state.confirmed_root_cause}*")

    action_suggestions = [
        "Review and update standard operating procedures (SOPs).",
        "Retrain operators on revised processes.",
        "Calibrate or replace faulty measurement equipment.",
        "Implement automated monitoring to flag anomalies in real time.",
        "Redesign component tolerance to accommodate variation.",
        "Apply control charting using our Statistical Software to monitor critical metrics.",
        "Conduct DOE (Design of Experiments) to optimize parameters.",
        "Verify root cause elimination through pilot run." 
    ]

    selected_actions = []
    for i, action in enumerate(action_suggestions):
        edit = st.text_area(f"Suggested Action {i+1}", value=action, key=f"ai_action_{i}")
        selected_actions.append(edit)

    st.markdown("You can now enter the final actions to implement, along with ownership and timeline:")

num_corrective = st.number_input("How many corrective actions would you like to plan?", min_value=1, max_value=10, step=1, key="num_corrective")
for i in range(num_corrective):
    st.markdown(f"#### Corrective Action {i+1}")
    description = st.text_area(f"Describe the action (Step {i+1})", key=f"corr_desc_{i}")
    responsible = st.text_input(f"Responsible Person/Team (Step {i+1})", key=f"corr_owner_{i}")
    start_date = st.date_input(f"Start Date (Step {i+1})", value=datetime.date.today(), key=f"corr_start_{i}")
    end_date = st.date_input(f"Planned Completion Date (Step {i+1})", value=datetime.date.today(), key=f"corr_end_{i}")
    expected_outcome = st.text_area(f"Expected Outcome (Step {i+1})", key=f"corr_outcome_{i}")
    corrective_actions.append({
        "description": description,
        "responsible": responsible,
        "start_date": start_date,
        "end_date": end_date,
        "expected_outcome": expected_outcome
    })

if st.button("💡 Save Corrective Action Plan"):
    st.success("Corrective actions recorded.")
    st.write("### 🧾 Permanent Corrective Action Plan Summary")
    for i, action in enumerate(corrective_actions):
        st.markdown(f"**{i+1}.** {action['description']}")
        st.markdown(f"&nbsp;&nbsp;👤 *Responsible:* {action['responsible']}")
        st.markdown(f"&nbsp;&nbsp;📆 *Start:* {action['start_date']} → *End:* {action['end_date']}")
        st.markdown(f"&nbsp;&nbsp;🎯 *Expected Outcome:* {action['expected_outcome']}")

st.caption("D5: Define and implement long-term actions to eliminate the root cause and prevent recurrence.")

# ---------------------------
# Section: D6 – Validation of Actions
# ---------------------------
st.markdown("---")
st.header("D6 – Validation of Actions")

validated_actions = []
num_validations = st.number_input("How many corrective actions will you validate?", min_value=1, max_value=10, step=1, key="num_validations")

for i in range(num_validations):
    st.markdown(f"#### Validation Entry {i+1}")
    validated_action = st.text_area(f"Which action is being validated? (Entry {i+1})", key=f"val_action_{i}")
    validation_method = st.text_input(f"Method of Validation (Entry {i+1})", key=f"val_method_{i}")
    validation_result = st.radio(
        f"Was the action effective? (Entry {i+1})",
        options=["Yes", "Partially", "No"],
        key=f"val_result_{i}"
    )
    notes = st.text_area(f"Notes or Observations (Entry {i+1})", key=f"val_notes_{i}")
    validated_actions.append({
        "validated_action": validated_action,
        "method": validation_method,
        "result": validation_result,
        "notes": notes
    })

if st.button("✅ Save Validation Results"):
    st.success("Validation results recorded.")
    st.write("### 📊 Corrective Action Validation Summary")
    for i, val in enumerate(validated_actions):
        st.markdown(f"**{i+1}.** *{val['validated_action']}*")
        st.markdown(f"&nbsp;&nbsp;🔬 *Method:* {val['method']}")
        st.markdown(f"&nbsp;&nbsp;📈 *Result:* {val['result']}")
        st.markdown(f"&nbsp;&nbsp;📝 *Notes:* {val['notes']}")

st.caption("D6: Confirm that the corrective actions were effective and did not cause unintended consequences.")

# ---------------------------
# Section: D7 – Prevent Recurrence
# ---------------------------
st.markdown("---")
st.header("D7 – Prevent Recurrence")

fmea_updates = st.text_area("📄 Describe updates made to the FMEA (Failure Mode and Effects Analysis):")
cp_updates = st.text_area("🗂️ Describe updates made to the Control Plan:")

additional_preventive = st.text_area("✅ List any other preventive measures taken to ensure the issue does not recur:")

if st.button("💾 Save Recurrence Prevention Details"):
    st.success("Recurrence prevention measures saved.")
    st.write("### 🛡️ Recurrence Prevention Summary")
    st.markdown(f"**FMEA Updates:** {fmea_updates}")
    st.markdown(f"**Control Plan Updates:** {cp_updates}")
    st.markdown(f"**Additional Measures:** {additional_preventive}")

st.caption("D7: Modify documentation and take systemic actions to ensure the issue does not happen again.")

# ---------------------------
# Section: D8 – Congratulate the Team
# ---------------------------
st.markdown("---")
st.header("D8 – Congratulate the Team")

team_success_message = st.text_area("🎉 Write a congratulatory message for the team:",
                                    value="Congratulations team! Your collaboration and efforts have resolved the issue successfully.")

generate_certificate = st.checkbox("🏅 Generate virtual team recognition certificate")
generate_pdf = st.checkbox("📄 Generate PDF summary")
generate_excel = st.checkbox("📊 Generate Excel summary")
generate_pptx = st.checkbox("📽️ Generate PowerPoint presentation")

if st.button("🎊 Finalize & Celebrate"):
    st.success("👏 Team recognition submitted!")
    st.write("### 🎉 Message to the Team")
    st.markdown(f"{team_success_message}")

    if generate_certificate:
        st.markdown("#### 🏆 Virtual Certificate")
        st.markdown("This certifies that the 8D Team has successfully completed a problem-solving process and contributed to continuous improvement.")
        st.markdown("*Issued by the Quality Department*")

    if generate_pdf:
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        pdf.multi_cell(0, 10, txt="8D Problem Solving Summary\n\n" + team_success_message)
        pdf.output("8D_Report.pdf")
        with open("8D_Report.pdf", "rb") as f:
            st.download_button("📥 Download PDF Report", f, file_name="8D_Report.pdf")

    if generate_excel:
        df = pd.DataFrame({
            "Discipline": [f"D{i+1}" for i in range(8)],
            "Description": [
                "Team Formation", "Problem Description", "Containment Actions", "Root Cause Analysis",
                "Corrective Actions", "Validation", "Prevent Recurrence", "Congratulate Team"
            ]
        })
        df.to_excel("8D_Report.xlsx", index=False)
        with open("8D_Report.xlsx", "rb") as f:
            st.download_button("📥 Download Excel Summary", f, file_name="8D_Report.xlsx")

    if generate_pptx:
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = "8D Problem Solving Report"
        slide.placeholders[1].text = "Final summary presentation"

        disciplines = [
            ("D1 – Team Formation", "Define the cross-functional team."),
            ("D2 – Problem Description", "Clearly describe the issue (5W2H + Fishbone)."),
            ("D3 – Containment Actions", "Short-term actions to contain the problem."),
            ("D4 – Root Cause Analysis", "5 Whys analysis to find the root cause."),
            ("D5 – Corrective Actions", "Permanent solutions to eliminate the root cause."),
            ("D6 – Validation", "Prove that corrective actions are effective."),
            ("D7 – Prevent Recurrence", "Update FMEA, control plans, and take preventive actions."),
            ("D8 – Congratulate Team", team_success_message)
        ]

        for title, content in disciplines:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = title
            slide.placeholders[1].text = content

        prs.save("8D_Presentation.pptx")
        with open("8D_Presentation.pptx", "rb") as f:
            st.download_button("📥 Download PowerPoint Presentation", f, file_name="8D_Presentation.pptx")

st.caption("D8: Recognize the team's efforts and document the success of the 8D process with export options.")
